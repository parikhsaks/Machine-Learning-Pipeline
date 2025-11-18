"""
Transformer Model Presentation Generator - Educational Version
Creates a detailed 6-slide presentation with comprehensive explanations
Designed for easy understanding and teaching
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_transformer_presentation():
    """Create an educational presentation explaining Transformer model architecture"""
    # Initialize presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Define color scheme
    PRIMARY_BLUE = RGBColor(41, 98, 255)    # Bright blue
    DARK_BLUE = RGBColor(10, 25, 55)        # Dark navy
    ACCENT_GREEN = RGBColor(34, 139, 34)    # Forest green for highlights
    ACCENT_ORANGE = RGBColor(255, 140, 0)   # Orange for important points
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 240)
    DARK_GRAY = RGBColor(64, 64, 64)
    
    # ============ SLIDE 1: INTRODUCTION & MOTIVATION ============
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Gradient background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE
    
    # Main title
    title_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(1), Inches(9), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "The Transformer Architecture"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle with key insight
    subtitle_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(2.3), Inches(9), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Revolutionizing NLP with Parallel Attention"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Key innovation box
    innovation_box = slide1.shapes.add_textbox(
        Inches(1.5), Inches(3.3), Inches(7), Inches(1.2)
    )
    innovation_frame = innovation_box.text_frame
    innovation_para = innovation_frame.paragraphs[0]
    innovation_para.text = "Key Innovation:"
    innovation_para.font.size = Pt(20)
    innovation_para.font.bold = True
    innovation_para.font.color.rgb = WHITE
    
    inn_p2 = innovation_frame.add_paragraph()
    inn_p2.text = "Replace sequential processing (RNN/LSTM) with"
    inn_p2.font.size = Pt(16)
    inn_p2.font.color.rgb = WHITE
    inn_p2.alignment = PP_ALIGN.CENTER
    
    inn_p3 = innovation_frame.add_paragraph()
    inn_p3.text = "parallel self-attention mechanism"
    inn_p3.font.size = Pt(16)
    inn_p3.font.color.rgb = WHITE
    inn_p3.alignment = PP_ALIGN.CENTER
    
    # Citation
    ref_box = slide1.shapes.add_textbox(
        Inches(0.5), Inches(4.8), Inches(9), Inches(0.5)
    )
    ref_frame = ref_box.text_frame
    ref_frame.text = "Vaswani et al., 2017 - Google Brain & Google Research"
    ref_para = ref_frame.paragraphs[0]
    ref_para.font.size = Pt(14)
    ref_para.font.italic = True
    ref_para.font.color.rgb = WHITE
    ref_para.alignment = PP_ALIGN.CENTER
    
    # ============ SLIDE 2: WHY TRANSFORMERS? THE PROBLEM ============
    slide_layout = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box2 = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.7)
    )
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Why Transformers? Understanding the Problem"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(34)
    title_para2.font.bold = True
    title_para2.font.color.rgb = DARK_BLUE
    
    # Problems with RNNs
    problem_box = slide2.shapes.add_textbox(
        Inches(0.5), Inches(1), Inches(4.3), Inches(4)
    )
    problem_frame = problem_box.text_frame
    problem_para = problem_frame.paragraphs[0]
    problem_para.text = "RNN/LSTM Limitations:"
    problem_para.font.size = Pt(18)
    problem_para.font.bold = True
    problem_para.font.color.rgb = RGBColor(220, 20, 60)  # Crimson
    
    prob_p1 = problem_frame.add_paragraph()
    prob_p1.text = "\n❌ Sequential Processing"
    prob_p1.font.size = Pt(15)
    prob_p1.font.bold = True
    
    prob_p2 = problem_frame.add_paragraph()
    prob_p2.text = "Must process word-by-word"
    prob_p2.font.size = Pt(13)
    prob_p2.font.color.rgb = DARK_GRAY
    
    prob_p3 = problem_frame.add_paragraph()
    prob_p3.text = "Can't parallelize = SLOW training"
    prob_p3.font.size = Pt(13)
    prob_p3.font.color.rgb = DARK_GRAY
    
    prob_p4 = problem_frame.add_paragraph()
    prob_p4.text = "\n❌ Long-Range Dependencies"
    prob_p4.font.size = Pt(15)
    prob_p4.font.bold = True
    
    prob_p5 = problem_frame.add_paragraph()
    prob_p5.text = "Information gets 'diluted'"
    prob_p5.font.size = Pt(13)
    prob_p5.font.color.rgb = DARK_GRAY
    
    prob_p6 = problem_frame.add_paragraph()
    prob_p6.text = "Gradient vanishing over distance"
    prob_p6.font.size = Pt(13)
    prob_p6.font.color.rgb = DARK_GRAY
    
    prob_p7 = problem_frame.add_paragraph()
    prob_p7.text = "\n❌ Fixed-Size Hidden State"
    prob_p7.font.size = Pt(15)
    prob_p7.font.bold = True
    
    prob_p8 = problem_frame.add_paragraph()
    prob_p8.text = "Bottleneck for information"
    prob_p8.font.size = Pt(13)
    prob_p8.font.color.rgb = DARK_GRAY
    
    # Transformer solutions
    solution_box = slide2.shapes.add_textbox(
        Inches(5.2), Inches(1), Inches(4.3), Inches(4)
    )
    solution_frame = solution_box.text_frame
    solution_para = solution_frame.paragraphs[0]
    solution_para.text = "Transformer Solutions:"
    solution_para.font.size = Pt(18)
    solution_para.font.bold = True
    solution_para.font.color.rgb = ACCENT_GREEN
    
    sol_p1 = solution_frame.add_paragraph()
    sol_p1.text = "\n✅ Parallel Processing"
    sol_p1.font.size = Pt(15)
    sol_p1.font.bold = True
    sol_p1.font.color.rgb = ACCENT_GREEN
    
    sol_p2 = solution_frame.add_paragraph()
    sol_p2.text = "All positions simultaneously"
    sol_p2.font.size = Pt(13)
    sol_p2.font.color.rgb = DARK_GRAY
    
    sol_p3 = solution_frame.add_paragraph()
    sol_p3.text = "10-100x faster training!"
    sol_p3.font.size = Pt(13)
    sol_p3.font.color.rgb = DARK_GRAY
    
    sol_p4 = solution_frame.add_paragraph()
    sol_p4.text = "\n✅ Direct Connections"
    sol_p4.font.size = Pt(15)
    sol_p4.font.bold = True
    sol_p4.font.color.rgb = ACCENT_GREEN
    
    sol_p5 = solution_frame.add_paragraph()
    sol_p5.text = "Any word sees all words"
    sol_p5.font.size = Pt(13)
    sol_p5.font.color.rgb = DARK_GRAY
    
    sol_p6 = solution_frame.add_paragraph()
    sol_p6.text = "No information decay"
    sol_p6.font.size = Pt(13)
    sol_p6.font.color.rgb = DARK_GRAY
    
    sol_p7 = solution_frame.add_paragraph()
    sol_p7.text = "\n✅ Dynamic Representation"
    sol_p7.font.size = Pt(15)
    sol_p7.font.bold = True
    sol_p7.font.color.rgb = ACCENT_GREEN
    
    sol_p8 = solution_frame.add_paragraph()
    sol_p8.text = "Attention weights adapt"
    sol_p8.font.size = Pt(13)
    sol_p8.font.color.rgb = DARK_GRAY
    
    # ============ SLIDE 3: SELF-ATTENTION INTUITION ============
    slide_layout = prs.slide_layouts[6]
    slide3 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box3 = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.7)
    )
    title_frame3 = title_box3.text_frame
    title_frame3.text = "Self-Attention: The Core Innovation"
    title_para3 = title_frame3.paragraphs[0]
    title_para3.font.size = Pt(34)
    title_para3.font.bold = True
    title_para3.font.color.rgb = DARK_BLUE
    
    # Intuitive explanation
    intuition_box = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.95), Inches(9), Inches(1)
    )
    intuition_frame = intuition_box.text_frame
    intuition_para = intuition_frame.paragraphs[0]
    intuition_para.text = "💡 Intuition: For each word, figure out which other words to 'pay attention to'"
    intuition_para.font.size = Pt(16)
    intuition_para.font.italic = True
    intuition_para.font.color.rgb = PRIMARY_BLUE
    intuition_para.alignment = PP_ALIGN.CENTER
    
    # Example sentence
    example_box = slide3.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(9), Inches(0.5)
    )
    example_frame = example_box.text_frame
    example_para = example_frame.paragraphs[0]
    example_para.text = 'Example: "The cat sat on the mat because it was tired"'
    example_para.font.size = Pt(15)
    example_para.font.bold = True
    example_para.alignment = PP_ALIGN.CENTER
    
    # Process explanation
    process_box = slide3.shapes.add_textbox(
        Inches(0.5), Inches(2.4), Inches(9), Inches(2.7)
    )
    process_frame = process_box.text_frame
    
    proc_p1 = process_frame.paragraphs[0]
    proc_p1.text = "How Self-Attention Works:"
    proc_p1.font.size = Pt(18)
    proc_p1.font.bold = True
    proc_p1.font.color.rgb = DARK_BLUE
    
    proc_p2 = process_frame.add_paragraph()
    proc_p2.text = "\n1️⃣ Create Three Representations for Each Word:"
    proc_p2.font.size = Pt(15)
    proc_p2.font.bold = True
    
    proc_p3 = process_frame.add_paragraph()
    proc_p3.text = "   • Query (Q): What am I looking for?"
    proc_p3.font.size = Pt(14)
    
    proc_p4 = process_frame.add_paragraph()
    proc_p4.text = '   • Key (K): What do I contain?'
    proc_p4.font.size = Pt(14)
    
    proc_p5 = process_frame.add_paragraph()
    proc_p5.text = '   • Value (V): What information do I provide?'
    proc_p5.font.size = Pt(14)
    
    proc_p6 = process_frame.add_paragraph()
    proc_p6.text = '\n2️⃣ Calculate Attention Scores:'
    proc_p6.font.size = Pt(15)
    proc_p6.font.bold = True
    
    proc_p7 = process_frame.add_paragraph()
    proc_p7.text = '   • "it" (Query) × all words (Keys) = compatibility scores'
    proc_p7.font.size = Pt(14)
    
    proc_p8 = process_frame.add_paragraph()
    proc_p8.text = '   • High score with "cat" → "it" refers to "cat"'
    proc_p8.font.size = Pt(14)
    
    proc_p9 = process_frame.add_paragraph()
    proc_p9.text = "\n3️⃣ Apply Softmax → Get attention weights (probabilities)"
    proc_p9.font.size = Pt(15)
    proc_p9.font.bold = True
    
    proc_p10 = process_frame.add_paragraph()
    proc_p10.text = "\n4️⃣ Weighted sum of Values → Context-aware representation!"
    proc_p10.font.size = Pt(15)
    proc_p10.font.bold = True
    
    # ============ SLIDE 4: ATTENTION MECHANICS & MULTI-HEAD ============
    slide_layout = prs.slide_layouts[6]
    slide4 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box4 = slide4.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.7)
    )
    title_frame4 = title_box4.text_frame
    title_frame4.text = "Attention Formula & Multi-Head Attention"
    title_para4 = title_frame4.paragraphs[0]
    title_para4.font.size = Pt(34)
    title_para4.font.bold = True
    title_para4.font.color.rgb = DARK_BLUE
    
    # Mathematical formula section
    math_box = slide4.shapes.add_textbox(
        Inches(0.5), Inches(0.95), Inches(9), Inches(1.8)
    )
    math_frame = math_box.text_frame
    
    math_p1 = math_frame.paragraphs[0]
    math_p1.text = "The Attention Formula:"
    math_p1.font.size = Pt(18)
    math_p1.font.bold = True
    math_p1.font.color.rgb = PRIMARY_BLUE
    
    math_p2 = math_frame.add_paragraph()
    math_p2.text = "\n         Attention(Q,K,V) = softmax(QK^T / √d_k) × V"
    math_p2.font.size = Pt(16)
    math_p2.font.name = "Courier New"
    math_p2.font.bold = True
    math_p2.alignment = PP_ALIGN.CENTER
    
    math_p3 = math_frame.add_paragraph()
    math_p3.text = "\n📌 Why divide by √d_k?"
    math_p3.font.size = Pt(15)
    math_p3.font.bold = True
    
    math_p4 = math_frame.add_paragraph()
    math_p4.text = "• Prevents dot products from getting too large (variance control)"
    math_p4.font.size = Pt(13)
    
    math_p5 = math_frame.add_paragraph()
    math_p5.text = "• Keeps gradients stable during training"
    math_p5.font.size = Pt(13)
    
    math_p6 = math_frame.add_paragraph()
    math_p6.text = "• Without it: softmax becomes too 'peaked' (one-hot)"
    math_p6.font.size = Pt(13)
    
    # Multi-head attention section
    multihead_box = slide4.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(9), Inches(2.2)
    )
    multihead_frame = multihead_box.text_frame
    
    mh_p1 = multihead_frame.paragraphs[0]
    mh_p1.text = "Multi-Head Attention: Different Perspectives"
    mh_p1.font.size = Pt(18)
    mh_p1.font.bold = True
    mh_p1.font.color.rgb = ACCENT_ORANGE
    
    mh_p2 = multihead_frame.add_paragraph()
    mh_p2.text = "\n🎯 Analogy: Like having 8 different 'experts' look at the sentence"
    mh_p2.font.size = Pt(14)
    mh_p2.font.italic = True
    
    mh_p3 = multihead_frame.add_paragraph()
    mh_p3.text = "\nEach head might learn to focus on:"
    mh_p3.font.size = Pt(14)
    mh_p3.font.bold = True
    
    mh_p4 = multihead_frame.add_paragraph()
    mh_p4.text = "• Head 1: Subject-verb relationships"
    mh_p4.font.size = Pt(13)
    
    mh_p5 = multihead_frame.add_paragraph()
    mh_p5.text = "• Head 2: Pronoun references"
    mh_p5.font.size = Pt(13)
    
    mh_p6 = multihead_frame.add_paragraph()
    mh_p6.text = "• Head 3: Adjacent word dependencies"
    mh_p6.font.size = Pt(13)
    
    mh_p7 = multihead_frame.add_paragraph()
    mh_p7.text = "• Head 4-8: Other linguistic patterns"
    mh_p7.font.size = Pt(13)
    
    mh_p8 = multihead_frame.add_paragraph()
    mh_p8.text = "\n→ Concatenate all heads → Linear projection → Rich representation!"
    mh_p8.font.size = Pt(14)
    mh_p8.font.bold = True
    
    # ============ SLIDE 5: COMPLETE ARCHITECTURE ============
    slide_layout = prs.slide_layouts[6]
    slide5 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box5 = slide5.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.7)
    )
    title_frame5 = title_box5.text_frame
    title_frame5.text = "Complete Transformer Architecture"
    title_para5 = title_frame5.paragraphs[0]
    title_para5.font.size = Pt(34)
    title_para5.font.bold = True
    title_para5.font.color.rgb = DARK_BLUE
    
    # Encoder section
    encoder_box = slide5.shapes.add_textbox(
        Inches(0.5), Inches(0.95), Inches(4.3), Inches(4.3)
    )
    encoder_frame = encoder_box.text_frame
    
    enc_p1 = encoder_frame.paragraphs[0]
    enc_p1.text = "📥 ENCODER (Left Side)"
    enc_p1.font.size = Pt(18)
    enc_p1.font.bold = True
    enc_p1.font.color.rgb = PRIMARY_BLUE
    
    enc_p2 = encoder_frame.add_paragraph()
    enc_p2.text = "\nInput Processing:"
    enc_p2.font.size = Pt(15)
    enc_p2.font.bold = True
    
    enc_p3 = encoder_frame.add_paragraph()
    enc_p3.text = "1. Word Embeddings (512-dim)"
    enc_p3.font.size = Pt(13)
    
    enc_p4 = encoder_frame.add_paragraph()
    enc_p4.text = "2. + Positional Encoding"
    enc_p4.font.size = Pt(13)
    
    enc_p5 = encoder_frame.add_paragraph()
    enc_p5.text = "   (sine/cosine patterns)"
    enc_p5.font.size = Pt(12)
    enc_p5.font.italic = True
    
    enc_p6 = encoder_frame.add_paragraph()
    enc_p6.text = "\n6 Identical Layers, each has:"
    enc_p6.font.size = Pt(15)
    enc_p6.font.bold = True
    
    enc_p7 = encoder_frame.add_paragraph()
    enc_p7.text = "• Multi-Head Attention"
    enc_p7.font.size = Pt(13)
    
    enc_p8 = encoder_frame.add_paragraph()
    enc_p8.text = "• Add & Normalize"
    enc_p8.font.size = Pt(13)
    
    enc_p9 = encoder_frame.add_paragraph()
    enc_p9.text = "• Feed-Forward Network"
    enc_p9.font.size = Pt(13)
    
    enc_p10 = encoder_frame.add_paragraph()
    enc_p10.text = "  (512→2048→512)"
    enc_p10.font.size = Pt(12)
    enc_p10.font.italic = True
    
    enc_p11 = encoder_frame.add_paragraph()
    enc_p11.text = "• Add & Normalize"
    enc_p11.font.size = Pt(13)
    
    enc_p12 = encoder_frame.add_paragraph()
    enc_p12.text = "\n💡 Processes entire input"
    enc_p12.font.size = Pt(14)
    enc_p12.font.bold = True
    
    enc_p13 = encoder_frame.add_paragraph()
    enc_p13.text = "   simultaneously!"
    enc_p13.font.size = Pt(14)
    enc_p13.font.bold = True
    
    # Decoder section
    decoder_box = slide5.shapes.add_textbox(
        Inches(5.2), Inches(0.95), Inches(4.3), Inches(4.3)
    )
    decoder_frame = decoder_box.text_frame
    
    dec_p1 = decoder_frame.paragraphs[0]
    dec_p1.text = "📤 DECODER (Right Side)"
    dec_p1.font.size = Pt(18)
    dec_p1.font.bold = True
    dec_p1.font.color.rgb = ACCENT_ORANGE
    
    dec_p2 = decoder_frame.add_paragraph()
    dec_p2.text = "\nOutput Generation:"
    dec_p2.font.size = Pt(15)
    dec_p2.font.bold = True
    
    dec_p3 = decoder_frame.add_paragraph()
    dec_p3.text = "1. Previous outputs"
    dec_p3.font.size = Pt(13)
    
    dec_p4 = decoder_frame.add_paragraph()
    dec_p4.text = "2. + Positional Encoding"
    dec_p4.font.size = Pt(13)
    
    dec_p5 = decoder_frame.add_paragraph()
    dec_p5.text = "\n6 Identical Layers, each has:"
    dec_p5.font.size = Pt(15)
    dec_p5.font.bold = True
    
    dec_p6 = decoder_frame.add_paragraph()
    dec_p6.text = "• Masked Self-Attention"
    dec_p6.font.size = Pt(13)
    
    dec_p7 = decoder_frame.add_paragraph()
    dec_p7.text = "  (can't see future!)"
    dec_p7.font.size = Pt(12)
    dec_p7.font.italic = True
    dec_p7.font.color.rgb = RGBColor(220, 20, 60)
    
    dec_p8 = decoder_frame.add_paragraph()
    dec_p8.text = "• Add & Normalize"
    dec_p8.font.size = Pt(13)
    
    dec_p9 = decoder_frame.add_paragraph()
    dec_p9.text = "• Cross-Attention"
    dec_p9.font.size = Pt(13)
    dec_p9.font.bold = True
    
    dec_p10 = decoder_frame.add_paragraph()
    dec_p10.text = "  (attend to encoder!)"
    dec_p10.font.size = Pt(12)
    dec_p10.font.italic = True
    dec_p10.font.color.rgb = ACCENT_GREEN
    
    dec_p11 = decoder_frame.add_paragraph()
    dec_p11.text = "• Add & Normalize"
    dec_p11.font.size = Pt(13)
    
    dec_p12 = decoder_frame.add_paragraph()
    dec_p12.text = "• Feed-Forward"
    dec_p12.font.size = Pt(13)
    
    dec_p13 = decoder_frame.add_paragraph()
    dec_p13.text = "• Add & Normalize"
    dec_p13.font.size = Pt(13)
    
    dec_p14 = decoder_frame.add_paragraph()
    dec_p14.text = "\n→ Linear → Softmax"
    dec_p14.font.size = Pt(14)
    dec_p14.font.bold = True
    
    dec_p15 = decoder_frame.add_paragraph()
    dec_p15.text = "→ Next word probability!"
    dec_p15.font.size = Pt(14)
    dec_p15.font.bold = True
    
    # ============ SLIDE 6: KEY INSIGHTS & IMPACT ============
    slide_layout = prs.slide_layouts[6]
    slide6 = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box6 = slide6.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.7)
    )
    title_frame6 = title_box6.text_frame
    title_frame6.text = "Key Insights & Revolutionary Impact"
    title_para6 = title_frame6.paragraphs[0]
    title_para6.font.size = Pt(34)
    title_para6.font.bold = True
    title_para6.font.color.rgb = DARK_BLUE
    
    # Key insights
    insights_box = slide6.shapes.add_textbox(
        Inches(0.5), Inches(0.95), Inches(4.3), Inches(4)
    )
    insights_frame = insights_box.text_frame
    
    ins_p1 = insights_frame.paragraphs[0]
    ins_p1.text = "🎯 Key Insights"
    ins_p1.font.size = Pt(18)
    ins_p1.font.bold = True
    ins_p1.font.color.rgb = PRIMARY_BLUE
    
    ins_p2 = insights_frame.add_paragraph()
    ins_p2.text = "\n1. Parallelization"
    ins_p2.font.size = Pt(15)
    ins_p2.font.bold = True
    
    ins_p3 = insights_frame.add_paragraph()
    ins_p3.text = "• No sequential bottleneck"
    ins_p3.font.size = Pt(13)
    
    ins_p4 = insights_frame.add_paragraph()
    ins_p4.text = "• GPU-friendly computation"
    ins_p4.font.size = Pt(13)
    
    ins_p5 = insights_frame.add_paragraph()
    ins_p5.text = "\n2. Attention Visualization"
    ins_p5.font.size = Pt(15)
    ins_p5.font.bold = True
    
    ins_p6 = insights_frame.add_paragraph()
    ins_p6.text = "• See what model focuses on"
    ins_p6.font.size = Pt(13)
    
    ins_p7 = insights_frame.add_paragraph()
    ins_p7.text = "• Interpretable decisions"
    ins_p7.font.size = Pt(13)
    
    ins_p8 = insights_frame.add_paragraph()
    ins_p8.text = "\n3. Transfer Learning"
    ins_p8.font.size = Pt(15)
    ins_p8.font.bold = True
    
    ins_p9 = insights_frame.add_paragraph()
    ins_p9.text = "• Pre-train on massive data"
    ins_p9.font.size = Pt(13)
    
    ins_p10 = insights_frame.add_paragraph()
    ins_p10.text = "• Fine-tune for any task"
    ins_p10.font.size = Pt(13)
    
    ins_p11 = insights_frame.add_paragraph()
    ins_p11.text = "\n4. Scalability"
    ins_p11.font.size = Pt(15)
    ins_p11.font.bold = True
    
    ins_p12 = insights_frame.add_paragraph()
    ins_p12.text = "• GPT-3: 175B parameters"
    ins_p12.font.size = Pt(13)
    
    ins_p13 = insights_frame.add_paragraph()
    ins_p13.text = "• Consistent improvements"
    ins_p13.font.size = Pt(13)
    
    # Impact section
    impact_box = slide6.shapes.add_textbox(
        Inches(5.2), Inches(0.95), Inches(4.3), Inches(4)
    )
    impact_frame = impact_box.text_frame
    
    imp_p1 = impact_frame.paragraphs[0]
    imp_p1.text = "🚀 Revolutionary Impact"
    imp_p1.font.size = Pt(18)
    imp_p1.font.bold = True
    imp_p1.font.color.rgb = ACCENT_ORANGE
    
    imp_p2 = impact_frame.add_paragraph()
    imp_p2.text = "\nNLP Breakthroughs:"
    imp_p2.font.size = Pt(15)
    imp_p2.font.bold = True
    
    imp_p3 = impact_frame.add_paragraph()
    imp_p3.text = "• BERT (Google, 2018)"
    imp_p3.font.size = Pt(13)
    
    imp_p4 = impact_frame.add_paragraph()
    imp_p4.text = "• GPT Series (OpenAI)"
    imp_p4.font.size = Pt(13)
    
    imp_p5 = impact_frame.add_paragraph()
    imp_p5.text = "• ChatGPT, Claude, etc."
    imp_p5.font.size = Pt(13)
    
    imp_p6 = impact_frame.add_paragraph()
    imp_p6.text = "\nBeyond Text:"
    imp_p6.font.size = Pt(15)
    imp_p6.font.bold = True
    
    imp_p7 = impact_frame.add_paragraph()
    imp_p7.text = "• Vision (ViT, DALL-E)"
    imp_p7.font.size = Pt(13)
    
    imp_p8 = impact_frame.add_paragraph()
    imp_p8.text = "• Audio (Whisper)"
    imp_p8.font.size = Pt(13)
    
    imp_p9 = impact_frame.add_paragraph()
    imp_p9.text = "• Protein Folding (AlphaFold)"
    imp_p9.font.size = Pt(13)
    
    imp_p10 = impact_frame.add_paragraph()
    imp_p10.text = "• Code (Copilot, Codex)"
    imp_p10.font.size = Pt(13)
    
    imp_p11 = impact_frame.add_paragraph()
    imp_p11.text = '\n"Attention is literally'
    imp_p11.font.size = Pt(14)
    imp_p11.font.italic = True
    imp_p11.font.bold = True
    
    imp_p12 = impact_frame.add_paragraph()
    imp_p12.text = 'all you need!"'
    imp_p12.font.size = Pt(14)
    imp_p12.font.italic = True
    imp_p12.font.bold = True
    
    # Footer tip
    footer_box = slide6.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(9), Inches(0.4)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "💡 Remember: Transformers = Parallel Attention + No Recurrence!"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.font.bold = True
    footer_para.font.color.rgb = DARK_BLUE
    footer_para.alignment = PP_ALIGN.CENTER
    
    return prs

def add_slide_numbers(prs):
    """Add slide numbers to all slides except the title slide"""
    GRAY = RGBColor(128, 128, 128)
    
    for i, slide in enumerate(prs.slides):
        if i == 0:  # Skip title slide
            continue
            
        # Add slide number
        slide_num_box = slide.shapes.add_textbox(
            Inches(9.2), Inches(5.2), Inches(0.5), Inches(0.3)
        )
        slide_num_frame = slide_num_box.text_frame
        slide_num_frame.text = f"{i + 1}/6"
        slide_num_para = slide_num_frame.paragraphs[0]
        slide_num_para.font.size = Pt(12)
        slide_num_para.font.color.rgb = GRAY
        slide_num_para.alignment = PP_ALIGN.CENTER

def save_presentation(prs, filename="transformer_explained.pptx"):
    """Save the presentation to a PPTX file"""
    # Add slide numbers
    add_slide_numbers(prs)
    
    # Save file
    prs.save(filename)
    print("=" * 60)
    print(f"✅ Presentation saved as '{filename}'")
    print("=" * 60)
    print("\n📊 SLIDE CONTENTS & TALKING POINTS:\n")
    
    print("SLIDE 1: Introduction & Motivation")
    print("-" * 40)
    print("• Start with the problem: 'Why did we need something new?'")
    print("• Emphasize the paradigm shift from sequential to parallel")
    print("• Mention this paper changed everything in AI")
    
    print("\nSLIDE 2: Why Transformers? The Problem")
    print("-" * 40)
    print("• Use hand gestures to show sequential vs parallel")
    print("• Give concrete example: translating a 100-word sentence")
    print("• RNN: 100 sequential steps, Transformer: 1 parallel step")
    print("• Emphasize the training speed improvement (10-100x)")
    
    print("\nSLIDE 3: Self-Attention Intuition")
    print("-" * 40)
    print("• Use the 'it' example - everyone understands pronouns")
    print("• Explain Q, K, V with analogy: dating app matching")
    print("  - Query: What you're looking for")
    print("  - Key: What others offer")
    print("  - Value: The actual person/information")
    print("• Show how 'it' connects to 'cat' through high attention")
    
    print("\nSLIDE 4: Attention Formula & Multi-Head")
    print("-" * 40)
    print("• Don't get too mathematical - focus on intuition")
    print("• Scaling: 'Like adjusting volume to prevent distortion'")
    print("• Multi-head: 'Different experts looking at different aspects'")
    print("• One head might focus on grammar, another on meaning")
    
    print("\nSLIDE 5: Complete Architecture")
    print("-" * 40)
    print("• Walk through an example: 'Hello' → 'Bonjour'")
    print("• Encoder understands 'Hello', Decoder generates 'Bonjour'")
    print("• Cross-attention: Where translation actually happens")
    print("• Masking: Can't cheat by looking at future words")
    
    print("\nSLIDE 6: Key Insights & Impact")
    print("-" * 40)
    print("• Connect to current AI boom (ChatGPT, etc.)")
    print("• Emphasize this is foundation of all modern LLMs")
    print("• Show enthusiasm about the elegance of the solution")
    print("• End with: 'Questions about any part?'")
    
    print("\n" + "=" * 60)
    print("📝 KEY QUESTIONS TO PREPARE FOR:")
    print("-" * 60)
    
    print("Q1: Why is it called 'self-attention'?")
    print("A: Because each word attends to ALL words in the same")
    print("   sequence, including itself.")
    
    print("\nQ2: How is this different from CNN?")
    print("A: CNNs have local receptive fields, Transformers see")
    print("   everything at once. CNNs need many layers for long-range.")
    
    print("\nQ3: What's the computational complexity?")
    print("A: O(n²·d) where n is sequence length, d is dimension.")
    print("   Quadratic in length but constant depth (vs RNN's O(n)).")
    
    print("\nQ4: Why positional encoding?")
    print("A: Without it, 'cat loves dog' = 'dog loves cat'")
    print("   (attention has no inherent order awareness).")
    
    print("\nQ5: Main limitation?")
    print("A: Quadratic memory for long sequences. That's why")
    print("   context windows are limited (GPT-4: 32k tokens).")
    
    print("\n" + "=" * 60)
    print("🎯 QUICK EXPLANATION TEMPLATE:")
    print("-" * 60)
    print('"The Transformer replaces sequential processing with')
    print('parallel attention. Instead of reading word-by-word')
    print('like RNNs, it processes all words simultaneously,')
    print('with each word deciding which other words to focus on.')
    print('This attention mechanism, combined with positional')
    print('encoding, captures both meaning and sequence order.')
    print('The result: 100x faster training and better performance!"')
    
    print("\n" + "=" * 60)
    print("💪 YOU'VE GOT THIS! The slides are clear and intuitive.")
    print("Remember: Enthusiasm is contagious - if you're excited")
    print("about how elegant Transformers are, your professor will be too!")
    print("=" * 60)
    
    return filename

# ============ MAIN EXECUTION ============
if __name__ == "__main__":
    print("\n" + "🤖 TRANSFORMER PRESENTATION GENERATOR" + "\n")
    print("Creating educational presentation with detailed explanations...")
    print("-" * 60)
    
    # Create the presentation
    presentation = create_transformer_presentation()
    
    # Save to file
    output_file = save_presentation(presentation)
    
    print("\n✨ Good luck with your presentation, Sakshi! 🎓")